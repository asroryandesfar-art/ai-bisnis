"""
Generate the 10-slide BotNesia Investor Pitch Deck (Investor-Readiness Assets,
Phase 19). Uses python-pptx DIRECTLY rather than document_generator.generate_pptx()
-- that shared helper hardcodes an auto-title-slide subtitle ("Dibuat otomatis
oleh BotNesia") meant for ad-hoc AI-generated documents, which isn't appropriate
for a polished investor cover slide. This script gives full control over the
cover slide and a consistent dark/brand-colored background across all 10 slides,
matching the same visual identity as the landing page (frontend/landing.css).

Run once: python3 scripts/generate_pitch_deck.py
Output: docs/marketing/BotNesia-Investor-Pitch-Deck.pptx
Source-of-truth content (editable without touching Python):
        docs/marketing/pitch-deck-content.md
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).parent.parent / "docs" / "marketing"

# Same palette as frontend/styles.css :root tokens, for visual consistency
# with the landing page / app.
BG = RGBColor(0x09, 0x0B, 0x10)
SURFACE = RGBColor(0x15, 0x1A, 0x24)
TEXT = RGBColor(0xF4, 0xF6, 0xF8)
TEXT_2 = RGBColor(0xA8, 0xB0, 0xBF)
BRAND = RGBColor(0x8B, 0x7C, 0xFF)
CYAN = RGBColor(0x51, 0xD7, 0xE8)
GREEN = RGBColor(0x45, 0xD3, 0x9B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)
CONTENT_W = SLIDE_W - 2 * MARGIN


def _set_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    return box


def _set_run(paragraph, text, size, color, bold=False, font="Manrope"):
    """Set text + formatting on a fresh, blank paragraph (every call site here
    only ever uses freshly created paragraphs, never appends a 2nd run)."""
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.name = font


def add_content_slide(prs, eyebrow, title, bullets=None, body=None, number=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)

    eyebrow_box = _textbox(slide, MARGIN, Inches(0.55), CONTENT_W, Inches(0.4))
    p = eyebrow_box.text_frame.paragraphs[0]
    _set_run(p, eyebrow.upper(), 13, CYAN, bold=True)

    title_box = _textbox(slide, MARGIN, Inches(1.0), CONTENT_W, Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    _set_run(p, title, 32, TEXT, bold=True)

    if bullets:
        body_box = _textbox(slide, MARGIN, Inches(2.15), CONTENT_W, Inches(4.6))
        tf = body_box.text_frame
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(14)
            _set_run(p, f"●  {bullet}", 17, TEXT_2)
    elif body:
        body_box = _textbox(slide, MARGIN, Inches(2.15), CONTENT_W, Inches(4.6))
        tf = body_box.text_frame
        for i, para in enumerate(body.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(12)
            _set_run(p, para, 17, TEXT_2)

    if number is not None:
        num_box = _textbox(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.55), Inches(0.6), Inches(0.4))
        p = num_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        _set_run(p, f"{number:02d} / 10", 10, TEXT_2)

    return slide


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1 — Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    logo_path = Path(__file__).parent.parent / "frontend" / "public" / "assets" / "brand" / "botnesia-clean-logo.png"
    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(5.92), Inches(1.6), height=Inches(1.5))
    title_box = _textbox(slide, MARGIN, Inches(3.5), CONTENT_W, Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p, "BotNesia", 54, TEXT, bold=True)
    sub_box = _textbox(slide, MARGIN, Inches(4.5), CONTENT_W, Inches(0.7))
    p = sub_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p, "AI Workforce Platform for Indonesian Businesses", 20, CYAN)

    # Slide 2 — Problem
    add_content_slide(
        prs, "Problem", "Masalah Bisnis Indonesia",
        bullets=[
            "Customer service mahal — gaji tim besar, training berulang",
            "Respon ke pelanggan lambat, apalagi di luar jam kerja",
            "Knowledge produk tercecer di banyak orang & dokumen",
            "Tim kecil, kewalahan menangani lonjakan chat & order",
            "Tidak ada analitik — keputusan bisnis berdasar tebakan",
        ], number=2,
    )

    # Slide 3 — Solution
    add_content_slide(
        prs, "Solution", "BotNesia sebagai AI Workforce",
        body=(
            "BotNesia menghadirkan tim AI lengkap — bukan satu chatbot, tapi satu\n"
            "tim AI Workforce: Customer Service, Sales, Marketing, Finance, HR,\n"
            "Operations, Security, hingga Executive Assistant — yang bekerja 24/7\n"
            "dan saling terhubung, tanpa bisnis perlu membangun tim teknologi sendiri."
        ), number=3,
    )

    # Slide 4 — Product
    add_content_slide(
        prs, "Product", "Semua Fitur Utama",
        bullets=[
            "Multi Agent System (Supervisor, CS, Sales, FAQ, Knowledge, Trainer, Memory)",
            "AI Workforce lintas-domain (Finance, Marketing, HR, Operations, Security, Executive)",
            "Executive Center & AI Business Analyst — root cause, rekomendasi, action plan otomatis",
            "Communication Center — WhatsApp, Instagram, Facebook, Telegram, Website, Email",
            "Knowledge Base AI — auto-generate FAQ & SOP dari dokumen",
            "Investor Demo Mode — AI menganalisis bisnis secara live, dalam hitungan detik",
        ], number=4,
    )

    # Slide 5 — Technology
    add_content_slide(
        prs, "Technology", "Multi Agent Architecture",
        bullets=[
            "Supervisor Agent mengoordinasikan seluruh agent spesialis",
            "FastAPI + PostgreSQL — backend multi-tenant, terisolasi penuh per workspace",
            "Groq LLM — inferensi cepat untuk seluruh kecerdasan AI agent",
            "Human-in-the-loop — setiap keputusan penting tetap memerlukan persetujuan manusia",
            "RBAC, audit log, dan enkripsi kredensial di setiap layer",
        ], number=5,
    )

    # Slide 6 — Market Opportunity
    add_content_slide(
        prs, "Market Opportunity", "Indonesia, UMKM, SME, Enterprise",
        bullets=[
            "Indonesia — salah satu ekonomi digital dengan pertumbuhan tercepat di Asia Tenggara",
            "60+ juta UMKM, sebagian besar belum terjangkau teknologi AI Workforce",
            "Bisnis menengah (SME) butuh otomasi CS & sales tanpa menambah tim besar",
            "Enterprise butuh AI Workforce multi-tenant dengan kendali akses korporat",
            "Satu platform, tiga segmen pasar — dari UMKM hingga enterprise",
        ], number=6,
    )

    # Slide 7 — Business Model
    add_content_slide(
        prs, "Business Model", "Subscription",
        bullets=[
            "Free — Rp0, 1 AI Agent, 100 percakapan/bulan",
            "Starter — Rp99rb/bulan, 2 AI Agents, Knowledge Base dasar",
            "Pro — Rp299rb/bulan, 5 AI Agents, WhatsApp integration, Analytics lengkap",
            "Business — Rp999rb/bulan, 10 AI Agents, Team Management, Priority Support",
            "Enterprise — Custom, Unlimited Agents, White Label, SLA, SSO",
        ], number=7,
    )

    # Slide 8 — Competitive Advantage
    add_content_slide(
        prs, "Competitive Advantage", "Mengapa BotNesia Berbeda",
        bullets=[
            "AI Workforce lengkap — bukan sekadar chatbot FAQ",
            "Human-in-the-loop di setiap keputusan penting, bukan AI yang berjalan sendiri",
            "Multi-tenant sejak awal — arsitektur siap untuk skala enterprise",
            "Dibangun untuk Indonesia — Bahasa Indonesia native, WhatsApp-first",
            "Executive Center menyatukan 6 domain bisnis jadi satu skor kesehatan bisnis",
        ], number=8,
    )

    # Slide 9 — Roadmap
    add_content_slide(
        prs, "Roadmap", "12 Bulan ke Depan",
        bullets=[
            "Q1 — Perluasan integrasi WhatsApp Business API & onboarding tenant baru",
            "Q2 — White-label & multi-tenant reseller untuk agency/partner",
            "Q3 — Perluasan AI Workforce Marketplace & template industri baru",
            "Q4 — Kemitraan strategis dengan inkubator & program AI nasional",
        ], number=9,
    )

    # Slide 10 — Vision
    add_content_slide(
        prs, "Vision", "Menjadi AI Workforce Platform Terbesar di Indonesia",
        body=(
            "Menjadi platform AI Workforce nomor satu di Indonesia — tempat setiap\n"
            "UMKM hingga perusahaan besar bisa memiliki tim AI selengkap perusahaan\n"
            "teknologi besar, tanpa harus membangun tim engineering sendiri."
        ), number=10,
    )

    return prs


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    out_path = OUT_DIR / "BotNesia-Investor-Pitch-Deck.pptx"
    prs.save(str(out_path))
    print(f"Wrote {out_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
