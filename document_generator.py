"""document_generator.py — Generate PDF / DOCX / XLSX / PPTX dari spec terstruktur.

Spec adalah dict sederhana (biasanya hasil satu panggilan `_call_llm_json`
yang mengubah permintaan user jadi outline):

    {
      "title": str,
      "sections": [{"heading": str, "body": str}, ...],
      "table_rows": [[str, ...], ...]  # opsional, header di baris pertama
      "slides": [{"title": str, "bullets": [str, ...]}, ...]  # opsional, khusus PPTX
      "logo_path": str  # opsional, path file gambar -- khusus PDF, ditaruh di atas judul
      "logo_width_inch": float  # opsional, default 1.0 -- lebar logo (tinggi proporsional)
    }

Setiap `generate_*` toleran terhadap field yang hilang — selalu balik bytes
yang valid, tidak pernah raise untuk spec yang "kurang lengkap" (hanya raise
untuk format yang tidak dikenal).
"""
from __future__ import annotations

import io
from xml.sax.saxutils import escape

SUPPORTED_FORMATS = {"pdf", "docx", "xlsx", "pptx"}


def normalize_spec(raw: dict | None, *, fallback_title: str = "Dokumen") -> dict:
    raw = raw if isinstance(raw, dict) else {}
    title = str(raw.get("title") or fallback_title).strip() or fallback_title

    sections = []
    for item in raw.get("sections") or []:
        if not isinstance(item, dict):
            continue
        heading = str(item.get("heading") or "").strip()
        body = str(item.get("body") or "").strip()
        if heading or body:
            sections.append({"heading": heading, "body": body})

    table_rows = []
    for row in raw.get("table_rows") or []:
        if isinstance(row, (list, tuple)):
            table_rows.append([str(cell) for cell in row])

    slides = []
    for item in raw.get("slides") or []:
        if not isinstance(item, dict):
            continue
        slide_title = str(item.get("title") or "").strip()
        bullets = [str(b).strip() for b in (item.get("bullets") or []) if str(b).strip()]
        if slide_title or bullets:
            slides.append({"title": slide_title, "bullets": bullets})

    logo_path = str(raw.get("logo_path") or "").strip() or None
    try:
        logo_width_inch = float(raw.get("logo_width_inch") or 1.0)
    except (TypeError, ValueError):
        logo_width_inch = 1.0

    return {
        "title": title, "sections": sections, "table_rows": table_rows, "slides": slides,
        "logo_path": logo_path, "logo_width_inch": logo_width_inch,
    }


def generate_pdf(spec: dict) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.lib.utils import ImageReader
    from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    spec = normalize_spec(spec)
    styles = getSampleStyleSheet()
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4)
    flow = []

    if spec["logo_path"]:
        try:
            reader = ImageReader(spec["logo_path"])
            iw, ih = reader.getSize()
            target_w = spec["logo_width_inch"] * inch
            flow.append(Image(spec["logo_path"], width=target_w, height=target_w * ih / iw))
            flow.append(Spacer(1, 10))
        except Exception:
            pass  # logo_path tidak valid/tidak terbaca -- jangan sampai gagalkan seluruh dokumen

    flow.append(Paragraph(escape(spec["title"]), styles["Title"]))
    flow.append(Spacer(1, 12))

    for section in spec["sections"]:
        if section["heading"]:
            flow.append(Paragraph(escape(section["heading"]), styles["Heading2"]))
        if section["body"]:
            for para in section["body"].split("\n"):
                if para.strip():
                    flow.append(Paragraph(escape(para.strip()), styles["BodyText"]))
        flow.append(Spacer(1, 8))

    if spec["table_rows"]:
        table = Table(spec["table_rows"])
        table.setStyle(TableStyle([
            ("GRID", (0, 0), (-1, -1), 0.5, "#999999"),
            ("BACKGROUND", (0, 0), (-1, 0), "#eeeeee"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
        ]))
        flow.append(table)

    if not spec["sections"] and not spec["table_rows"]:
        flow.append(Paragraph("Tidak ada konten.", styles["BodyText"]))

    doc.build(flow)
    return buf.getvalue()


def generate_docx(spec: dict) -> bytes:
    from docx import Document

    spec = normalize_spec(spec)
    document = Document()
    document.add_heading(spec["title"], level=0)

    for section in spec["sections"]:
        if section["heading"]:
            document.add_heading(section["heading"], level=2)
        if section["body"]:
            for para in section["body"].split("\n"):
                if para.strip():
                    document.add_paragraph(para.strip())

    if spec["table_rows"]:
        rows = spec["table_rows"]
        n_cols = max((len(r) for r in rows), default=0)
        if n_cols:
            table = document.add_table(rows=0, cols=n_cols)
            table.style = "Light Grid Accent 1"
            for row in rows:
                cells = table.add_row().cells
                for i in range(n_cols):
                    cells[i].text = row[i] if i < len(row) else ""

    if not spec["sections"] and not spec["table_rows"]:
        document.add_paragraph("Tidak ada konten.")

    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def generate_xlsx(spec: dict) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    spec = normalize_spec(spec)
    wb = Workbook()
    ws = wb.active
    ws.title = (spec["title"] or "Sheet1")[:31] or "Sheet1"

    row_idx = 1
    ws.cell(row=row_idx, column=1, value=spec["title"]).font = Font(bold=True, size=14)
    row_idx += 2

    if spec["table_rows"]:
        for row in spec["table_rows"]:
            for col_idx, value in enumerate(row, start=1):
                ws.cell(row=row_idx, column=col_idx, value=value)
            row_idx += 1
    elif spec["sections"]:
        ws.cell(row=row_idx, column=1, value="Bagian").font = Font(bold=True)
        ws.cell(row=row_idx, column=2, value="Isi").font = Font(bold=True)
        row_idx += 1
        for section in spec["sections"]:
            ws.cell(row=row_idx, column=1, value=section["heading"])
            ws.cell(row=row_idx, column=2, value=section["body"])
            row_idx += 1
    else:
        ws.cell(row=row_idx, column=1, value="Tidak ada konten.")

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def generate_pptx(spec: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    spec = normalize_spec(spec)
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = spec["title"]
    if title_slide.placeholders and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Dibuat otomatis oleh BotNesia"

    slide_specs = spec["slides"] or [
        {"title": s["heading"] or "Slide", "bullets": [s["body"]] if s["body"] else []}
        for s in spec["sections"]
    ]
    bullet_layout = prs.slide_layouts[1]
    for slide_spec in slide_specs:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = slide_spec.get("title") or "Slide"
        body = slide.placeholders[1].text_frame
        bullets = slide_spec.get("bullets") or []
        if not bullets:
            body.text = ""
        else:
            body.text = bullets[0]
            for bullet in bullets[1:]:
                p = body.add_paragraph()
                p.text = bullet

    if not slide_specs:
        empty = prs.slides.add_slide(bullet_layout)
        empty.shapes.title.text = "Tidak ada konten"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


_GENERATORS = {
    "pdf": generate_pdf,
    "docx": generate_docx,
    "xlsx": generate_xlsx,
    "pptx": generate_pptx,
}


def generate_document(fmt: str, spec: dict) -> tuple[bytes, str]:
    """Returns (bytes, content_type) for the requested format."""
    fmt = (fmt or "").strip().lower()
    if fmt not in _GENERATORS:
        raise ValueError(f"Format dokumen '{fmt}' tidak didukung. Pilih: {sorted(SUPPORTED_FORMATS)}")
    content_types = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    return _GENERATORS[fmt](spec), content_types[fmt]
